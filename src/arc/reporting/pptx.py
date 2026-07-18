"""Weekly PPTX slide generator using python-pptx (Tech Spec §9.1)."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from arc.paths import REPORTS_DIR

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

ARC_PRIMARY = RGBColor(204, 120, 92)
ARC_DARK = RGBColor(20, 20, 19)
ARC_MUTED = RGBColor(108, 106, 100)


def _set_bg(slide, color: RGBColor) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def write_weekly_pptx(week_str: str, context: dict[str, Any]) -> Path:
    """Generate a weekly PPTX deck from report context."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, ARC_DARK)
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(11), Inches(1.5)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = context.get("brand", "ARC")
    p.font.size = Pt(48)
    p.font.color.rgb = ARC_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = f"Weekly Report — {week_str}"
    p2.font.size = Pt(24)
    p2.font.color.rgb = ARC_MUTED

    # Slide 2: Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox2 = slide2.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(11), Inches(6)
    )
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    p.text = "本周概览"
    p.font.size = Pt(36)
    p.font.bold = True

    for h in context.get("headlines", []):
        p2 = tf2.add_paragraph()
        p2.text = f"  - {h.get('title', '')}"
        p2.font.size = Pt(18)
    if not context.get("headlines"):
        p2 = tf2.add_paragraph()
        p2.text = "  本周无显著新信号"
        p2.font.size = Pt(18)

    # Slide 3: Actions
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    txBox3 = slide3.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(11), Inches(6)
    )
    tf3 = txBox3.text_frame
    p = tf3.paragraphs[0]
    p.text = "下周任务"
    p.font.size = Pt(36)
    p.font.bold = True

    for i, a in enumerate(context.get("actions", []), 1):
        p2 = tf3.add_paragraph()
        p2.text = f"  {i}. {a}"
        p2.font.size = Pt(18)

    out_dir = REPORTS_DIR / "weekly"
    out_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = out_dir / f"{week_str}.pptx"
    prs.save(str(pptx_path))
    return pptx_path
